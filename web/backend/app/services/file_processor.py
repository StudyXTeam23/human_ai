"""File processing service for document upload."""
import base64
import os
from pathlib import Path
from typing import Tuple
import logging

from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation

logger = logging.getLogger(__name__)

# 上传文件保存目录
UPLOAD_DIR = Path(__file__).parent.parent.parent / "uploads"
UPLOAD_DIR.mkdir(exist_ok=True)


class FileProcessor:
    """Service for processing uploaded documents."""

    SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".pptx", ".txt"}
    MAX_FILE_SIZE = 40 * 1024 * 1024  # 40MB

    @staticmethod
    def validate_file(filename: str, file_size: int) -> None:
        """
        Validate uploaded file.

        Args:
            filename: Name of the uploaded file
            file_size: Size of the file in bytes

        Raises:
            ValueError: If file validation fails
        """
        # Check file size
        if file_size > FileProcessor.MAX_FILE_SIZE:
            raise ValueError(
                f"File size ({file_size / 1024 / 1024:.2f}MB) exceeds maximum allowed size (40MB)"
            )

        # Check file extension
        ext = Path(filename).suffix.lower()
        if ext not in FileProcessor.SUPPORTED_EXTENSIONS:
            raise ValueError(
                f"Unsupported file type: {ext}. Supported types: {', '.join(FileProcessor.SUPPORTED_EXTENSIONS)}"
            )

    @staticmethod
    def save_file(filename: str, content: bytes) -> str:
        """
        Save uploaded file to local storage.

        Args:
            filename: Name of the file
            content: File content as bytes

        Returns:
            str: Path to saved file
        """
        # Generate unique filename with timestamp
        import time
        timestamp = int(time.time() * 1000)
        safe_filename = f"{timestamp}_{filename}"
        file_path = UPLOAD_DIR / safe_filename

        # Save file
        with open(file_path, "wb") as f:
            f.write(content)

        logger.info(f"File saved: {file_path}")
        return str(file_path)

    @staticmethod
    def extract_text_from_pdf(file_path: str) -> str:
        """
        Extract text from PDF file.

        Args:
            file_path: Path to PDF file

        Returns:
            str: Extracted text
        """
        try:
            reader = PdfReader(file_path)
            text = ""
            for page in reader.pages:
                text += page.extract_text() + "\n"
            return text.strip()
        except Exception as e:
            logger.error(f"Error extracting text from PDF: {e}")
            raise ValueError(f"Failed to extract text from PDF: {str(e)}")

    @staticmethod
    def extract_text_from_docx(file_path: str) -> str:
        """
        Extract text from DOCX file.

        Args:
            file_path: Path to DOCX file

        Returns:
            str: Extracted text
        """
        try:
            doc = Document(file_path)
            text = "\n".join([paragraph.text for paragraph in doc.paragraphs])
            return text.strip()
        except Exception as e:
            logger.error(f"Error extracting text from DOCX: {e}")
            raise ValueError(f"Failed to extract text from DOCX: {str(e)}")

    @staticmethod
    def extract_text_from_pptx(file_path: str) -> str:
        """
        Extract text from PPTX file.

        Args:
            file_path: Path to PPTX file

        Returns:
            str: Extracted text
        """
        try:
            prs = Presentation(file_path)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text.strip()
        except Exception as e:
            logger.error(f"Error extracting text from PPTX: {e}")
            raise ValueError(f"Failed to extract text from PPTX: {str(e)}")

    @staticmethod
    def extract_text_from_txt(file_path: str) -> str:
        """
        Extract text from TXT file.

        Args:
            file_path: Path to TXT file

        Returns:
            str: Extracted text
        """
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                return f.read().strip()
        except UnicodeDecodeError:
            # Try with different encoding
            try:
                with open(file_path, "r", encoding="gbk") as f:
                    return f.read().strip()
            except Exception as e:
                logger.error(f"Error reading TXT file: {e}")
                raise ValueError(f"Failed to read TXT file: {str(e)}")
        except Exception as e:
            logger.error(f"Error reading TXT file: {e}")
            raise ValueError(f"Failed to read TXT file: {str(e)}")

    @staticmethod
    def process_file(file_path: str) -> Tuple[str, str]:
        """
        Process uploaded file and extract text.

        Args:
            file_path: Path to the uploaded file

        Returns:
            Tuple[str, str]: (extracted_text, base64_encoded_file)

        Raises:
            ValueError: If file processing fails
        """
        ext = Path(file_path).suffix.lower()

        # Extract text based on file type
        if ext == ".pdf":
            text = FileProcessor.extract_text_from_pdf(file_path)
        elif ext == ".docx":
            text = FileProcessor.extract_text_from_docx(file_path)
        elif ext == ".pptx":
            text = FileProcessor.extract_text_from_pptx(file_path)
        elif ext == ".txt":
            text = FileProcessor.extract_text_from_txt(file_path)
        else:
            raise ValueError(f"Unsupported file type: {ext}")

        # For document mode, we don't enforce length limits
        # The text will be sent directly to OpenAI
        logger.info(f"Extracted {len(text)} characters from document")

        # Convert file to base64
        with open(file_path, "rb") as f:
            file_content = f.read()
            base64_content = base64.b64encode(file_content).decode("utf-8")

        logger.info(
            f"File processed successfully: {len(text)} characters extracted"
        )

        return text, base64_content

    @staticmethod
    def cleanup_old_files(max_age_hours: int = 24) -> None:
        """
        Clean up old uploaded files.

        Args:
            max_age_hours: Maximum age of files to keep in hours
        """
        import time

        current_time = time.time()
        max_age_seconds = max_age_hours * 3600

        for file_path in UPLOAD_DIR.iterdir():
            if file_path.is_file():
                file_age = current_time - file_path.stat().st_mtime
                if file_age > max_age_seconds:
                    try:
                        file_path.unlink()
                        logger.info(f"Deleted old file: {file_path}")
                    except Exception as e:
                        logger.error(f"Failed to delete file {file_path}: {e}")

