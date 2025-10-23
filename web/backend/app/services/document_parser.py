"""Document parsing service for various file formats."""
import io
from typing import BinaryIO

import PyPDF2
from docx import Document
from pptx import Presentation


class DocumentParserService:
    """Service for parsing documents and extracting text."""

    def parse_pdf(self, file: BinaryIO) -> str:
        """
        Parse PDF file and extract text.

        Args:
            file: Binary file object

        Returns:
            Extracted text content
        """
        try:
            reader = PyPDF2.PdfReader(file)
            text = ""
            for page in reader.pages:
                text += page.extract_text() + "\n"
            return text.strip()
        except Exception as e:
            raise ValueError(f"Failed to parse PDF: {str(e)}")

    def parse_docx(self, file: BinaryIO) -> str:
        """
        Parse DOCX file and extract text.

        Args:
            file: Binary file object

        Returns:
            Extracted text content
        """
        try:
            doc = Document(file)
            text = "\n".join([para.text for para in doc.paragraphs if para.text.strip()])
            return text.strip()
        except Exception as e:
            raise ValueError(f"Failed to parse DOCX: {str(e)}")

    def parse_pptx(self, file: BinaryIO) -> str:
        """
        Parse PPTX file and extract text.

        Args:
            file: Binary file object

        Returns:
            Extracted text content
        """
        try:
            prs = Presentation(file)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text += shape.text + "\n"
            return text.strip()
        except Exception as e:
            raise ValueError(f"Failed to parse PPTX: {str(e)}")

    def parse_txt(self, file: BinaryIO) -> str:
        """
        Parse TXT file and extract text.

        Args:
            file: Binary file object

        Returns:
            Extracted text content
        """
        try:
            content = file.read()
            # Try UTF-8 first, fall back to latin-1
            try:
                return content.decode("utf-8").strip()
            except UnicodeDecodeError:
                return content.decode("latin-1").strip()
        except Exception as e:
            raise ValueError(f"Failed to parse TXT: {str(e)}")

    def parse(self, file: BinaryIO, filename: str) -> str:
        """
        Parse document based on file extension.

        Args:
            file: Binary file object
            filename: Original filename with extension

        Returns:
            Extracted text content

        Raises:
            ValueError: If file type is unsupported or parsing fails
        """
        ext = filename.lower().split(".")[-1]

        parsers = {
            "pdf": self.parse_pdf,
            "docx": self.parse_docx,
            "doc": self.parse_docx,  # Try DOCX parser for DOC
            "pptx": self.parse_pptx,
            "ppt": self.parse_pptx,  # Try PPTX parser for PPT
            "txt": self.parse_txt,
        }

        parser = parsers.get(ext)
        if not parser:
            raise ValueError(f"Unsupported file type: {ext}")

        return parser(file)

